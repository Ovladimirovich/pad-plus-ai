import logging
import re
from typing import List, Optional
from dataclasses import dataclass

logger = logging.getLogger("padplus.document_processor")

CHUNK_SIZE = 800
CHUNK_OVERLAP = 80
EMBEDDING_DIM = 1536


@dataclass
class DocumentChunk:
    index: int
    content: str


# ============================================================================
# TEXT EXTRACTION
# ============================================================================

def extract_text_from_txt(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


def extract_text_from_pdf(content: bytes) -> str:
    try:
        import fitz
    except ImportError:
        raise ImportError("PyMuPDF не установлен. Установите: pip install PyMuPDF")

    doc = fitz.open(stream=content, filetype="pdf")
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n\n".join(pages)


def extract_text_from_docx(content: bytes) -> str:
    try:
        import docx
    except ImportError:
        raise ImportError("python-docx не установлен. Установите: pip install python-docx")

    import io
    doc = docx.Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs)


def extract_text_from_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx не установлен. Установите: pip install python-pptx")

    import io
    prs = Presentation(io.BytesIO(content))
    slides = []
    for slide in prs.slides:
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        slides.append("\n".join(slide_texts))
    return "\n\n".join(slides)


def extract_text_from_md(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


EXTRACTORS = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".doc": extract_text_from_docx,
    ".pptx": extract_text_from_pptx,
    ".txt": extract_text_from_txt,
    ".md": extract_text_from_md,
    ".csv": extract_text_from_txt,
    ".json": extract_text_from_txt,
    ".xml": extract_text_from_txt,
    ".html": extract_text_from_txt,
}


def get_extractor(file_extension: str):
    ext = file_extension.lower()
    if ext not in EXTRACTORS:
        raise ValueError(f"Неподдерживаемый формат: {file_extension}")
    return EXTRACTORS[ext]


# ============================================================================
# CHUNKING
# ============================================================================

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[DocumentChunk]:
    if not text.strip():
        return []

    text = re.sub(r"\s+", " ", text).strip()
    words = text.split()
    chunks = []
    start = 0

    while start < len(words):
        end = start + chunk_size
        chunk_words = words[start:end]
        chunks.append(DocumentChunk(
            index=len(chunks),
            content=" ".join(chunk_words),
        ))
        step = chunk_size - overlap
        if step <= 0:
            break
        start += step

    return chunks


# ============================================================================
# MAIN PROCESSING PIPELINE
# ============================================================================

async def process_document(document_id: str, document_content: Optional[bytes] = None) -> bool:
    """
    Полный цикл обработки документа:
    1. Получить метаданные из БД
    2. Скачать файл из Supabase Storage
    3. Извлечь текст
    4. Разбить на чанки
    5. Получить эмбеддинги
    6. Сохранить в document_chunks
    7. Обновить статус документа

    Args:
        document_id: UUID документа
        document_content: Опционально — содержимое файла (если уже загружено в память)

    Returns:
        True если обработка успешна, иначе False
    """
    from core.supabase_client import get_supabase
    from core.pg_pool import get_connection, put_connection

    try:
        supabase = get_supabase()
        if not supabase:
            logger.error("Supabase не подключён")
            await _update_document_status(document_id, "failed")
            return False

        # Получаем метаданные документа
        doc_result = supabase.table("documents").select("*").eq("id", document_id).execute()
        if not doc_result.data:
            logger.error(f"Документ {document_id} не найден")
            return False

        doc = doc_result.data[0]
        user_id = doc.get("user_id")
        file_path = doc.get("file_path")
        filename = doc.get("filename", "")
        file_ext = "." + filename.split(".")[-1].lower() if "." in filename else ""

        # Обновляем статус на processing
        await _update_document_status(document_id, "processing")

        # Скачиваем файл, если не передан
        if document_content is None:
            try:
                storage_content = supabase.storage.from_("documents").download(file_path)
                document_content = storage_content
            except Exception as e:
                logger.error(f"Не удалось скачать {file_path}: {e}")
                await _update_document_status(document_id, "failed")
                return False

        # Извлекаем текст
        try:
            extractor = get_extractor(file_ext)
            text = extractor(document_content)
        except (ImportError, ValueError) as e:
            logger.error(f"Ошибка извлечения текста: {e}")
            await _update_document_status(document_id, "failed")
            return False

        if not text.strip():
            logger.warning(f"Документ {document_id} не содержит текста")
            await _update_document_status(document_id, "failed")
            return False

        # Разбиваем на чанки
        chunks = chunk_text(text, CHUNK_SIZE, CHUNK_OVERLAP)
        if not chunks:
            logger.warning(f"Документ {document_id}: нет чанков после разбиения")
            await _update_document_status(document_id, "failed")
            return False

        logger.info(f"Документ {document_id}: {len(chunks)} чанков, {len(text)} символов")

        # Получаем эмбеддинги
        embeddings = await _get_embeddings_for_chunks(chunks, user_id)
        if embeddings is None:
            logger.error(f"Не удалось получить эмбеддинги для {document_id}")
            await _update_document_status(document_id, "failed")
            return False

        # Сохраняем чанки в БД через прямой PostgreSQL
        conn = None
        try:
            conn = get_connection()
            cur = conn.cursor()

            # Удаляем старые чанки (если документ переобрабатывается)
            cur.execute("DELETE FROM document_chunks WHERE document_id = %s", (document_id,))

            for chunk, embedding in zip(chunks, embeddings):
                embedding_list = list(embedding)
                cur.execute(
                    """
                    INSERT INTO document_chunks (document_id, chunk_index, content, embedding, metadata)
                    VALUES (%s, %s, %s, %s::vector, %s)
                    """,
                    (
                        document_id,
                        chunk.index,
                        chunk.content,
                        embedding_list,
                        "{}",
                    ),
                )

            conn.commit()
            logger.info(f"Документ {document_id}: сохранено {len(chunks)} чанков с эмбеддингами")
        except Exception as e:
            logger.error(f"Ошибка сохранения чанков в БД: {e}")
            if conn:
                conn.rollback()
            await _update_document_status(document_id, "failed")
            return False
        finally:
            if conn:
                put_connection(conn)

        # Обновляем статус на completed
        await _update_document_status(document_id, "completed")

        # Обновляем summary в таблице documents (первые 500 символов)
        try:
            summary = text[:500].strip()
            supabase.table("documents").update({
                "summary": summary,
                "updated_at": "now()",
            }).eq("id", document_id).execute()
        except Exception as e:
            logger.warning(f"Не удалось обновить summary документа: {e}")

        return True

    except Exception as e:
        logger.error(f"Критическая ошибка обработки документа {document_id}: {e}", exc_info=True)
        await _update_document_status(document_id, "failed")
        return False


async def _update_document_status(document_id: str, status: str):
    """Обновляет статус документа в Supabase."""
    try:
        from core.supabase_client import get_supabase
        supabase = get_supabase()
        supabase.table("documents").update({
            "status": status,
        }).eq("id", document_id).execute()
    except Exception as e:
        logger.warning(f"Не удалось обновить статус {document_id} -> {status}: {e}")


async def _get_embeddings_for_chunks(chunks: List[DocumentChunk], user_id: str) -> Optional[List[List[float]]]:
    """Получает эмбеддинги для списка чанков через OpenRouter."""
    from runtime.llm_service import get_llm_service

    try:
        llm = get_llm_service()

        # Ищем API ключ OpenRouter пользователя
        api_key = await _get_user_openrouter_key(user_id)
        if not api_key:
            logger.warning(f"Нет OpenRouter ключа у пользователя {user_id}, пробую дефолтный")
            api_key = llm.default_api_key

        if not api_key:
            logger.error("Нет API ключа для получения эмбеддингов")
            return None

        texts = [chunk.content for chunk in chunks]

        # OpenRouter имеет лимит на размер batch — делим по 20 текстов
        all_embeddings = []
        batch_size = 20
        for i in range(0, len(texts), batch_size):
            batch = texts[i:i + batch_size]
            batch_embeddings = await llm.get_embeddings(batch, api_key=api_key)
            all_embeddings.extend(batch_embeddings)

        return all_embeddings

    except Exception as e:
        logger.error(f"Ошибка получения эмбеддингов: {e}")
        return None


async def _get_user_openrouter_key(user_id: str) -> Optional[str]:
    """Получает OpenRouter API ключ пользователя из таблицы user_api_keys."""
    try:
        from core.supabase_client import get_supabase
        from core.encryption import get_encryptor
        supabase = get_supabase()
        result = supabase.table("user_api_keys").select("api_key_encrypted").eq("user_id", user_id).eq("provider", "openrouter").execute()
        if result.data:
            encrypted = result.data[0].get("api_key_encrypted")
            if encrypted:
                encryptor = get_encryptor()
                return encryptor.decrypt(encrypted)
        return None
    except Exception as e:
        logger.warning(f"Не удалось получить ключ OpenRouter: {e}")
        return None


# ============================================================================
# SEARCH
# ============================================================================

async def search_document_chunks(
    query: str,
    user_id: str,
    limit: int = 5,
    similarity_threshold: float = 0.5,
) -> List[dict]:
    """
    Ищет релевантные чанки документов по векторному сходству.

    1. Получает эмбеддинг запроса
    2. Выполняет pgvector cosine similarity search
    3. Возвращает чанки документов пользователя

    Args:
        query: Поисковый запрос
        user_id: ID пользователя
        limit: Максимум результатов
        similarity_threshold: Минимальное сходство (0.0 - 1.0)

    Returns:
        Список словарей с полями: content, document_id, chunk_index, similarity, document_title
    """
    from core.pg_pool import get_connection, put_connection

    conn = None
    try:
        conn = get_connection()
        cur = conn.cursor()

        # Получаем эмбеддинг запроса
        from runtime.llm_service import get_llm_service
        llm = get_llm_service()
        api_key = await _get_user_openrouter_key(user_id) or llm.default_api_key
        if not api_key:
            logger.warning("Нет API ключа для поиска по документам")
            return []

        try:
            query_embedding = await llm.get_embeddings([query], api_key=api_key)
        except Exception as e:
            logger.error(f"Ошибка получения эмбеддинга запроса: {e}")
            return []

        if not query_embedding:
            return []

        query_vec = list(query_embedding[0])

        # Векторный поиск через pgvector
        cur.execute(
            """
            SELECT
                dc.id,
                dc.content,
                dc.document_id,
                dc.chunk_index,
                1 - (dc.embedding <=> %s::vector) AS similarity,
                d.title AS document_title
            FROM document_chunks dc
            JOIN documents d ON d.id = dc.document_id
            WHERE d.user_id = %s
              AND d.is_deleted = False
              AND 1 - (dc.embedding <=> %s::vector) > %s
            ORDER BY similarity DESC
            LIMIT %s
            """,
            (query_vec, user_id, query_vec, similarity_threshold, limit),
        )

        results = []
        for row in cur.fetchall():
            results.append({
                "chunk_id": str(row[0]),
                "content": row[1],
                "document_id": str(row[2]),
                "chunk_index": row[3],
                "similarity": round(float(row[4]), 4),
                "document_title": row[5],
            })

        return results

    except Exception as e:
        logger.error(f"Ошибка поиска по документам: {e}")
        return []
    finally:
        if conn:
            put_connection(conn)
